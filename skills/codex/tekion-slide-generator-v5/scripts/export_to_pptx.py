#!/usr/bin/env python3
"""
スライド画像をPowerPoint (PPTX)にエクスポートするスクリプト

指定ディレクトリ内の画像ファイルを番号順に読み込み、
16:9アスペクト比でPowerPointプレゼンテーションとして出力します。
各画像はスライド全体にぴったりフィットされます。

使用方法:
    python export_to_pptx.py --input-dir <画像ディレクトリ> --output <出力PPTXパス>
"""
import os
import sys
import argparse
import re
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


VERSION_PATTERN = re.compile(r'^(.+_\d+)(?:_v(\d+))?\.png$')


def natural_sort_key(s):
    """
    自然順ソート用のキー生成関数

    例: slide_1.png, slide_2.png, slide_10.png を正しく順序付け
    """
    return [int(text) if text.isdigit() else text.lower()
            for text in re.split('([0-9]+)', str(s))]


def group_and_select_latest_versions(image_files):
    """
    画像ファイルをベース名でグループ化し、各グループから最新バージョンを選択する

    バージョニング規則:
        ryoko-news-pipeline_01.png      → ベース名: ryoko-news-pipeline_01, v1（暗黙）
        ryoko-news-pipeline_01_v2.png   → ベース名: ryoko-news-pipeline_01, v2
        ryoko-news-pipeline_01_v3.png   → ベース名: ryoko-news-pipeline_01, v3

    バージョンなしファイルのみの場合は動作変化なし（完全後方互換）。

    Args:
        image_files: Pathオブジェクトのリスト

    Returns:
        最新バージョンのみのPathオブジェクトリスト（natural_sort_keyでソート済み）
    """
    groups = {}

    for img_path in image_files:
        match = VERSION_PATTERN.match(img_path.name)
        if match:
            base_name = match.group(1)
            version = int(match.group(2)) if match.group(2) else 1
        else:
            base_name = img_path.stem
            version = 1

        if base_name not in groups or version > groups[base_name][1]:
            groups[base_name] = (img_path, version)

    selected = [path for path, _ in groups.values()]
    # course_title（表紙）を常に先頭にソート
    return sorted(selected, key=lambda p: (0 if 'course_title' in p.name else 1, natural_sort_key(p)))


def export_to_pptx(input_dir: str, output_path: str, pattern: str = "*.png") -> bool:
    """
    画像ファイルをPowerPoint (PPTX)にエクスポート

    Args:
        input_dir: 入力画像ディレクトリ
        output_path: 出力PPTXファイルパス
        pattern: ファイルパターン（デフォルト: *.png）

    Returns:
        bool: 成功時True、失敗時False
    """
    try:
        input_path = Path(input_dir)

        if not input_path.exists():
            print(f"❌ エラー: 入力ディレクトリが存在しません: {input_dir}", file=sys.stderr)
            return False

        # 画像ファイルを取得（バージョン管理: 各スライドの最新バージョンのみ選択）
        all_images = list(input_path.glob(pattern))
        image_files = group_and_select_latest_versions(all_images)

        if not image_files:
            print(f"❌ エラー: 画像ファイルが見つかりません（パターン: {pattern}）", file=sys.stderr)
            return False

        print(f"📄 {len(image_files)}枚の画像を検出しました", file=sys.stderr)

        # PowerPointプレゼンテーションを作成
        prs = Presentation()
        # 16:9レイアウトに設定（デフォルトは4:3）
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        # スライドサイズ（インチ単位）
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        slide_aspect_ratio = slide_width / slide_height

        # 各画像をスライドとして追加（PNG 自体にフッターが焼き込まれているため、そのまま埋め込む）
        for img_file in image_files:
            try:
                # 画像を開いてサイズを取得
                with Image.open(img_file) as img:
                    img_width, img_height = img.size

                print(f"  ✓ {img_file.name}: {img_width}x{img_height}", file=sys.stderr)

                # 空白スライドを追加
                blank_slide_layout = prs.slide_layouts[6]  # 空白レイアウト
                slide = prs.slides.add_slide(blank_slide_layout)

                # 画像のアスペクト比を計算
                image_aspect_ratio = img_width / img_height

                # スライド全体にフィットさせる（アスペクト比を維持）
                if image_aspect_ratio > slide_aspect_ratio:
                    # 画像が横長: 幅をスライド幅に合わせる
                    w = slide_width
                    h = slide_width / image_aspect_ratio
                    x = Inches(0)
                    y = (slide_height - h) / 2
                else:
                    # 画像が縦長: 高さをスライド高さに合わせる
                    h = slide_height
                    w = slide_height * image_aspect_ratio
                    x = (slide_width - w) / 2
                    y = Inches(0)

                # 画像を追加
                slide.shapes.add_picture(
                    str(img_file),
                    x, y, w, h
                )

            except Exception as e:
                print(f"  ⚠️  {img_file.name}: 処理失敗 ({e})", file=sys.stderr)

        # 出力ディレクトリが存在しない場合は作成
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        # PowerPointファイルとして保存
        print(f"💾 PowerPoint生成中: {output_path}", file=sys.stderr)
        prs.save(output_path)

        file_size = os.path.getsize(output_path)
        print(f"✅ PowerPoint生成成功: {output_path} ({file_size:,} bytes, {len(prs.slides)}スライド)", file=sys.stderr)

        return True

    except Exception as e:
        print(f"❌ エラー: {type(e).__name__}: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc(file=sys.stderr)
        return False


def main():
    parser = argparse.ArgumentParser(
        description='スライド画像をPowerPoint (PPTX)にエクスポート',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
例:
  python export_to_pptx.py \\
    --input-dir ./slides_output/images \\
    --output ./slides_output/presentation.pptx
        """
    )
    parser.add_argument(
        '--input-dir',
        required=True,
        help='入力画像ディレクトリ（例: slides_output/images）'
    )
    parser.add_argument(
        '--output',
        required=True,
        help='出力PPTXファイルパス（例: slides_output/第四回.pptx）'
    )
    parser.add_argument(
        '--pattern',
        default='*.png',
        help='画像ファイルパターン（デフォルト: *.png）'
    )

    args = parser.parse_args()

    # PowerPoint生成
    success = export_to_pptx(args.input_dir, args.output, args.pattern)
    sys.exit(0 if success else 1)


if __name__ == '__main__':
    main()
